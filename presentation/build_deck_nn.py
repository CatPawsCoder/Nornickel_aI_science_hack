# -*- coding: utf-8 -*-
"""Презентация на официальном шаблоне «Норникель 2026».

Заполняет плейсхолдеры шаблона контентом проекта, сохраняя фирменный стиль:
Proxima Nova, фиолетовый 5302E0, мятный 00FFBF, серый 434343.
"""
import copy
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TPL = r"C:\Users\Asus ROG\Downloads\Норникель 2026.pptx"
OUT = r"C:\Users\Asus ROG\nauchny-klubok\presentation\nauchny-klubok-NN-v2.pptx"

VIOLET = RGBColor(0x53, 0x02, 0xE0)
MINT = RGBColor(0x00, 0xFF, 0xBF)
GREY = RGBColor(0x43, 0x43, 0x43)
DARK = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Proxima Nova"

prs = Presentation(TPL)
S = list(prs.slides)


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


# картинки-заглушки шаблона на слайдах 3 и 4 заслоняют текст — убираем
from pptx.enum.shapes import MSO_SHAPE_TYPE
for slide_idx in (2, 3):
    for sh in list(S[slide_idx].shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            delete_shape(sh)


def set_text(shape, lines, size=18, color=GREY, bold=False, spacing=1.15):
    """Полностью переписывает text_frame шейпа (список строк или (текст, опции))."""
    tf = shape.text_frame
    tf.word_wrap = True
    # очистить все параграфы кроме первого
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    items = lines if isinstance(lines, list) else [lines]
    for i, item in enumerate(items):
        text, opts = (item, {}) if isinstance(item, str) else item
        p = first if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = opts.get("color", color)


def add_box(slide, x, y, w, h, lines, size=18, color=GREY, bold=False, spacing=1.15,
            align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = lines if isinstance(lines, list) else [lines]
    for i, item in enumerate(items):
        text, opts = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", bold)
        run.font.color.rgb = opts.get("color", color)
    return box


# ============ S1 — ТИТУЛ (оставляем шапку хакатона, добавляем проект) ============
s = S[0]
add_box(s, 2.0, 4.35, 7.5, 0.55,
        [("НАУЧНЫЙ КЛУБОК", {"size": 30, "bold": True, "color": VIOLET})])
add_box(s, 2.0, 4.9, 7.6, 0.6,
        [("Верифицируемая карта знаний R&D для горно-металлургической отрасли",
          {"size": 15, "color": DARK})])

# ============ S2 — ПРОБЛЕМА ============
s = S[1]
set_text(s.shapes[0], [("ПРОБЛЕМА", {"size": 40, "bold": True, "color": DARK})])
set_text(s.shapes[1], [("Знания есть — общей структуры нет", {"size": 18, "bold": True, "color": VIOLET})])
set_text(s.shapes[2], [
    "• Институциональная память рассеяна по отчётам, презентациям и личным архивам",
    "• Команды дублируют литобзоры, не видя выполненных работ",
    "• Междисциплинарный поиск «режим + материал + условия» вручную занимает недели",
    "• Противоречивые выводы без единой верифицированной базы",
], size=15, color=GREY)

# ============ S3 — ОБРАЗ РЕШЕНИЯ ============
s = S[2]
set_text(s.shapes[0], [("ОБРАЗ РЕШЕНИЯ", {"size": 40, "bold": True, "color": DARK})])
set_text(s.shapes[1], [("Единый граф знаний вместо разрозненных документов",
                        {"size": 18, "bold": True, "color": VIOLET})])
set_text(s.shapes[2], [
    "Связываем 8 типов сущностей: публикации, эксперименты,",
    "технологии, материалы, оборудование, условия, выводы, эксперты.",
    "",
    "Запрос естественным языком: «методы обессоливания при сульфатах",
    "200–300 мг/л и сухом остатке ≤1000 мг/дм³» → структурированный",
    "ответ с источниками, уровнем достоверности и цитатами.",
], size=14, color=GREY)

# ============ S4 — АРХИТЕКТУРА ============
s = S[3]
set_text(s.shapes[0], [("АРХИТЕКТУРА: ГИБРИД", {"size": 40, "bold": True, "color": DARK})])
set_text(s.shapes[1], [("Граф — источник истины. LLM — интерфейс, но не источник фактов",
                        {"size": 18, "bold": True, "color": VIOLET})])
set_text(s.shapes[2], [
    "• Числа и диапазоны — только детерминированная regex-грамматика",
    "   (ошибки в концентрациях недопустимы — LLM их не генерирует)",
    "• Каждая цитата верифицируется string-match к исходному тексту:",
    "   галлюцинация физически не может попасть в граф",
    "• Сущности/связи — LLM по строгой JSON-схеме + тезаурус RU/EN",
    "   («электроэкстракция» = «electrowinning» = «ЭЭ»)",
    "• Хранение — Kùzu (Cypher), ответы — GraphRAG с цитатами",
], size=13.5, color=GREY)

# ============ S5 — МАСШТАБ ============
s = S[4]
# синие плашки-бейджи шаблона накладываются на заголовок — убираем их
delete_shape(s.shapes[5])
delete_shape(s.shapes[4])
title_box = s.shapes[1]
title_box.width = Inches(8.6)   # бейджей справа больше нет — заголовку можно всю ширину
title_box.height = Inches(0.5)
title_box.top = Inches(1.35)
set_text(title_box, [("МАСШТАБ НА ПОЛНОМ КОРПУСЕ", {"size": 24, "bold": True, "color": DARK})])
# подзаголовок ниже заголовка с запасом (шрифт-замена в PowerPoint шире растровой)
sub_box = s.shapes[2]
sub_box.top = Inches(2.05)
set_text(sub_box, [("4.7 ГБ архива обработано полностью, без пропусков",
                    {"size": 15, "bold": True, "color": VIOLET})])
# список — по одному пункту на строку, содержательнее про модель и граф
body_box = s.shapes[3]
body_box.top = Inches(2.6)
body_box.height = Inches(2.4)
set_text(body_box, [
    "• 1 828 уникальных документов (дедуп по контент-хэшу): PDF, DOCX, PPTX, XLS + OCR",
    "• 259 304 числовых условия — каждое дословно найдено в источнике (string-match)",
    "• 666 утверждений: строгая JSON-схема, дословная цитата хранится в узле графа",
    "• Строгий AND (вложенность интервалов): напр. 55–70°С + 150–300 А/м² → 3 из 1828",
    "• Граф Kùzu: 8 типов узлов, 19 типов рёбер, обход 3–4 уровней связей < 0.2 c",
    "• 3 000 рёбер «источники противоречат» + автопоиск пробелов в знаниях",
    "• Тезаурус RU/EN: 88 канонов, 400+ синонимов («электроэкстракция» = «electrowinning»)",
    "• Ответ на многопараметрический запрос — 2–3 секунды",
], size=12, color=GREY, spacing=1.2)

# ============ S6 — ДЕМО (две карточки) ============
s = S[5]
add_box(s, 0.4, 0.55, 9.0, 0.5, [("ДЕМО: ЭТАЛОННЫЕ ЗАПРОСЫ ТЗ", {"size": 30, "bold": True, "color": DARK})])
card1, card2 = s.shapes[0], s.shapes[1]
set_text(card1, [
    ("Обессоливание воды", {"size": 16, "bold": True, "color": VIOLET}),
    ("сульфаты/хлориды 200–300 мг/л, сухой остаток ≤1000 мг/дм³",
     {"size": 12, "color": GREY}),
    ("→ 23 верифицированных условия, методы HiPRO (98% воды), SULFATEQ (<300 мг/л)",
     {"size": 12, "color": GREY}),
    ("", {"size": 8}),
    ("Циркуляция католита (EW никеля)", {"size": 16, "bold": True, "color": VIOLET}),
    ("мировая практика, оптимальная скорость потока", {"size": 12, "color": GREY}),
    ("→ 60 утверждений: 20–30 л/ч через диафрагму, 1.5–5 м³/ч рециркуляция",
     {"size": 12, "color": GREY}),
])
set_text(card2, [
    ("Au, Ag, МПГ: штейн vs шлак", {"size": 16, "bold": True, "color": VIOLET}),
    ("эксперименты и публикации за 10 лет", {"size": 12, "color": GREY}),
    ("→ коэффициенты распределения: Au 1500, Pt 5000, Rh 7000–8000 при 65% Cu",
     {"size": 12, "color": GREY}),
    ("", {"size": 8}),
    ("Закачка шахтных вод", {"size": 16, "bold": True, "color": VIOLET}),
    ("Россия и зарубежье, ТЭП", {"size": 12, "color": GREY}),
    ("→ честный ответ: пробел в корпусе + внешний мониторинг (Yandex Search API)",
     {"size": 12, "color": GREY}),
])

# ============ S7 — ВЕРИФИКАЦИЯ ============
s = S[6]
add_box(s, 0.4, 0.55, 9.0, 0.5, [("МОДЕЛЬ ВЕРИФИКАЦИИ ЗНАНИЙ", {"size": 30, "bold": True, "color": DARK})])
add_box(s, 0.4, 1.25, 8.8, 0.4, [("Каждый факт: источник + дословная цитата + достоверность + дата",
                                  {"size": 16, "bold": True, "color": VIOLET})])
add_box(s, 0.4, 1.95, 9.0, 3.4, [
    "• 259 304 числа: каждое дословно найдено в источнике (string-match), 0 исключений",
    "• 666 LLM-утверждений: цитата не совпала → в граф не попадает; прошла — хранится в узле",
    "• Валидаторы ответа: числа, вещества и отрицания сверяются с фактами — иначе",
    "   ответ отдаётся дословными фактами (extractive), а не пересказом LLM",
    "• Уровни достоверности + версионирование: active / superseded с указанием замены",
    "• Зоны разногласий: 3 000 рёбер «источники расходятся» (параметр + вещество)",
    "• Пробелы: intent-якоря отличают отсутствующую тему от смежной (⚠️ в ответе)",
    "• Прослеживаемость: каждое число прямого ответа → ссылка [узел] → строка документа",
    "• Безопасность: санитизация вывода, аудит-лог по админ-токену, SHA256 данных",
], size=13, color=GREY)

# ============ S8 — ДЕПЛОЙ И МАСШТАБИРУЕМОСТЬ ============
s = S[7]
add_box(s, 0.3, 1.2, 8.3, 0.9, [("РАЗВЁРНУТО И МАСШТАБИРУЕМО", {"size": 34, "bold": True, "color": DARK})])
add_box(s, 0.3, 2.0, 8.3, 0.5, [("Живое демо: catpawws-ai-science-hack.hf.space",
                                 {"size": 17, "bold": True, "color": VIOLET})])
add_box(s, 0.3, 2.7, 9.2, 2.6, [
    "• Развёртывание одной командой: docker compose up, данные подтягиваются автоматически",
    "• Лайт-режим: SQLite FTS5 — рантайм < 1 ГБ RAM, без GPU (GPU нужен только сборке)",
    "• Новый домен = новые записи тезауруса, онтология не меняется",
    "• 3 LLM-провайдера с автопереключением: YandexGPT → GigaChat → DeepSeek (open-weights)",
    "• FAIR: каждый факт находим, доступен по REST, сериализуем в RDF, переиспользуем",
], size=14, color=GREY)

# ============ S9 — ИТОГ ============
s = S[8]
add_box(s, 0.4, 0.7, 9.0, 0.7, [("ГРАФ — ИСТОЧНИК ИСТИНЫ.", {"size": 32, "bold": True, "color": DARK})])
add_box(s, 0.4, 1.4, 9.0, 0.7, [("LLM — ИНТЕРФЕЙС, НО НЕ ФАКТ.", {"size": 32, "bold": True, "color": VIOLET})])
add_box(s, 0.4, 2.5, 9.0, 1.6, [
    "Готовая карта знаний R&D: многопараметрические запросы, верифицируемые",
    "ответы за секунды, границы знаний и противоречия — видимы.",
    "",
    "Roadmap: семантические эмбеддинги, цепочки процессов «выход→вход»,",
    "автомониторинг новых публикаций, дашборды руководителя.",
], size=15, color=GREY)
add_box(s, 0.4, 4.5, 9.0, 0.5, [
    ("Демо: catpawws-ai-science-hack.hf.space   •   Код: см. заявку",
     {"size": 14, "bold": True, "color": VIOLET})])

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
