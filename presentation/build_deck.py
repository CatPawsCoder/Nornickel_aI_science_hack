# -*- coding: utf-8 -*-
"""Генератор презентации «Научный клубок» на python-pptx (Node.js недоступен в среде)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- палитра «Руда и пламя» (Ore & Ember) ----
GRAPHITE = RGBColor(0x1B, 0x1F, 0x27)     # тёмный фон (титул/финал)
SLATE = RGBColor(0x3E, 0x54, 0x63)        # вторичный
STEEL = RGBColor(0x5B, 0x7B, 0x8C)        # поддерживающий
EMBER = RGBColor(0xE0, 0x8E, 0x45)        # акцент — расплавленная медь
EMBER_DARK = RGBColor(0xB8, 0x6A, 0x2A)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x26, 0x2E)
MUTED = RGBColor(0x6B, 0x74, 0x7D)
GREEN = RGBColor(0x3D, 0xB8, 0x7A)
CARD_BG = RGBColor(0xEC, 0xEE, 0xF0)

FONT_HEAD = "Cambria"
FONT_BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide(bg=LIGHT_BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def txt(slide, x, y, w, h, text, size=16, color=INK, bold=False, font=FONT_BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.bold = bold; run.font.name = font; run.font.italic = italic
    return box


def bullets(slide, x, y, w, h, items, size=14, color=INK, font=FONT_BODY,
            space_after=8, bullet_color=EMBER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        bu = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2])})
        buClr.append(srgb)
        pPr.set('indent', '-182880'); pPr.set('marL', '182880')
        pPr.append(buClr); pPr.append(buFont); pPr.append(bu)
        run = p.add_run(); run.text = item
        run.font.size = Pt(size); run.font.color.rgb = color; run.font.name = font
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=None, shadow=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def circle_num(slide, x, y, d, text, fill=EMBER, fg=WHITE, size=20):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = fg; r.font.name = FONT_HEAD
    return c


def page_num(slide, n):
    txt(slide, 12.5, 7.05, 0.6, 0.35, str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def kicker(slide, text, color=EMBER, dark=False):
    txt(slide, 0.7, 0.55, 8, 0.4, text.upper(), size=13, bold=True, color=color,
        font=FONT_BODY)


# ============================================================ SLIDE 1 — TITLE
s = add_slide(GRAPHITE)
# декоративные окружности (мотив — узел графа)
for (cx, cy, r, col, op) in [(11.6, 1.1, 1.6, SLATE, None), (12.4, 5.6, 1.1, EMBER_DARK, None)]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2*r), Inches(2*r))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.fill.transparency = 0
    c.line.fill.background(); c.shadow.inherit = False
    c.fill.fore_color.brightness = 0
sp = s.shapes[-1]
txt(s, 0.9, 1.5, 4, 0.5, "НАУЧНЫЙ ХАКАТОН · ГОРНАЯ МЕТАЛЛУРГИЯ", size=14, bold=True, color=EMBER, font=FONT_BODY)
txt(s, 0.85, 2.05, 10.5, 1.9, "Научный клубок", size=64, bold=True, color=WHITE, font=FONT_HEAD)
txt(s, 0.9, 3.35, 9.5, 1.0, "Верифицируемая карта знаний R&D для горно-металлургической отрасли",
    size=24, color=RGBColor(0xCF, 0xD6, 0xDC), font=FONT_BODY)
txt(s, 0.9, 4.35, 9.5, 0.6,
    "Граф знаний, где числа не может исказить ни одна модель — они проверяются "
    "дословным совпадением с источником.", size=15, italic=True, color=STEEL, font=FONT_BODY)
# нижняя строка с фактами
facts = ["101 публикация", "20 843 условия", "233 утверждения", "100% верификация"]
fx = 0.9
for f in facts:
    w = 0.16 * len(f) + 0.5
    rect(s, fx, 6.35, w, 0.55, fill=RGBColor(0x28, 0x2E, 0x39), radius=0.5)
    txt(s, fx, 6.35, w, 0.55, f, size=13, bold=True, color=EMBER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fx += w + 0.25
page_num(s, 1)

# ============================================================ SLIDE 2 — PROBLEM
s = add_slide(LIGHT_BG)
kicker(s, "Проблема")
txt(s, 0.7, 0.95, 10, 1.15, "Проблема не в отсутствии моделей —\nа в отсутствии общей структуры знаний",
    size=30, bold=True, color=INK, font=FONT_HEAD, line_spacing=1.05)
problems = [
    ("Потеря институциональной памяти", "Знания о выщелачивании, электролите, распределении металлов рассеяны по отчётам и личным архивам"),
    ("Дублирование усилий", "Команды заново пишут литобзоры, не видя, что аналогичная работа уже выполнена"),
    ("Сложность междисциплинарного поиска", "Нельзя быстро связать «кучное выщелачивание в холодном климате» с «выходом металла»"),
    ("Низкая скорость решений", "Ответ на вопрос о циркуляции электролита требует ручного сбора данных из десятков источников"),
]
cols = 2
cw, ch, gx, gy = 5.55, 1.85, 0.7, 2.35
for i, (title, body) in enumerate(problems):
    cx = gx + (i % cols) * (cw + 0.5)
    cy = gy + (i // cols) * (ch + 0.35)
    rect(s, cx, cy, cw, ch, fill=WHITE, shadow=True, radius=0.06)
    circle_num(s, cx + 0.3, cy + 0.3, 0.5, str(i+1), fill=EMBER, size=18)
    txt(s, cx + 1.0, cy + 0.28, cw - 1.3, 0.5, title, size=15, bold=True, color=INK, font=FONT_BODY)
    txt(s, cx + 1.0, cy + 0.82, cw - 1.3, ch - 1.0, body, size=12, color=MUTED, font=FONT_BODY, line_spacing=1.15)
page_num(s, 2)

# ============================================================ SLIDE 3 — SOLUTION IMAGE
s = add_slide(LIGHT_BG)
kicker(s, "Образ решения")
txt(s, 0.7, 0.95, 11.5, 1.15, "Единая карта знаний, связывающая\nразнородные сущности отрасли", size=30, bold=True,
    color=INK, font=FONT_HEAD, line_spacing=1.05)
entities = [
    ("Публикации", "Литобзоры, патенты, диссертации"),
    ("Эксперименты", "Протоколы, параметры процессов"),
    ("Технологии", "Схемы циркуляции, конструкции ячеек"),
    ("Материалы", "Сульфаты, никель, гипс, штейн"),
    ("Оборудование", "Ванны EW, печи ПВП, газоочистка"),
    ("Эксперты", "Лаборатории, авторы, компетенции"),
    ("Выводы", "Подтверждённые эффекты, ограничения"),
]
cw2, ch2 = 3.55, 1.35
positions = [(0.7,2.35),(4.5,2.35),(8.3,2.35),(0.7,3.95),(4.5,3.95),(8.3,3.95),(4.5,5.55)]
colors = [STEEL, SLATE, EMBER, STEEL, SLATE, EMBER, GRAPHITE]
for (title, body), (cx, cy), col in zip(entities, positions, colors):
    rect(s, cx, cy, cw2, ch2, fill=col, radius=0.08)
    txt(s, cx+0.25, cy+0.18, cw2-0.5, 0.4, title, size=15, bold=True, color=WHITE, font=FONT_BODY)
    txt(s, cx+0.25, cy+0.62, cw2-0.5, ch2-0.7, body, size=11, color=RGBColor(0xE4,0xE8,0xEB), font=FONT_BODY, line_spacing=1.1)
page_num(s, 3)

# ============================================================ SLIDE 4 — ARCHITECTURE (hybrid)
s = add_slide(GRAPHITE)
kicker(s, "Ключевая идея архитектуры", color=EMBER)
txt(s, 0.7, 0.95, 11.5, 0.7, "Гибрид с жёстким разделением ответственности", size=30, bold=True, color=WHITE, font=FONT_HEAD)
txt(s, 0.7, 1.55, 11, 0.5, "Граф — источник истины. LLM — интерфейс к нему на входе и выходе, но никогда не источник фактов.",
    size=15, italic=True, color=EMBER)

rows = [
    ("Числа, единицы, диапазоны", "Только детерминированная regex-грамматика", "ТЗ: ошибки в концентрациях недопустимы"),
    ("Верификация", "String-match каждой цитаты к исходному тексту", "Галлюцинация не может попасть в граф"),
    ("Сущности и связи", "LLM по строгой JSON-схеме + тезаурус-канонизатор", "RU/EN синонимы: «электроэкстракция» = «electrowinning»"),
    ("Хранение", "Kùzu — графовая СУБД с Cypher", "Обход 3-4 уровней связей < 0.2 с"),
    ("Синтез ответа", "GraphRAG: граф фильтрует, LLM оформляет", "Каждое утверждение = узел графа с цитатой"),
]
ty = 2.25
rh = 0.92
for i, (a, b, c) in enumerate(rows):
    cy = ty + i * (rh + 0.06)
    fill = RGBColor(0x25, 0x2B, 0x35) if i % 2 == 0 else RGBColor(0x1F, 0x24, 0x2D)
    rect(s, 0.7, cy, 11.9, rh, fill=fill, radius=0.05)
    txt(s, 0.95, cy + 0.12, 2.6, rh-0.2, a, size=13, bold=True, color=EMBER, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.7, cy + 0.12, 4.6, rh-0.2, b, size=12.5, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    txt(s, 8.4, cy + 0.12, 4.0, rh-0.2, c, size=11.5, italic=True, color=STEEL, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
page_num(s, 4)

# ============================================================ SLIDE 5 — PIPELINE
s = add_slide(LIGHT_BG)
kicker(s, "Как это работает")
txt(s, 0.7, 0.95, 11, 0.7, "Пайплайн: от документа до верифицированного ответа", size=28, bold=True, color=INK, font=FONT_HEAD)

steps = [
    ("Ингест", "PDF/DOCX RU/EN\n→ текст + метаданные"),
    ("Числа", "regex-грамматика\n20 843 условия"),
    ("Сущности", "LLM + тезаурус\n88 канонов, 400+ синонимов"),
    ("Граф", "Kùzu: 8 типов узлов\n19 типов рёбер"),
    ("Ответ", "BM25 + фильтры графа\n+ LLM-синтез с цитатами"),
]
n = len(steps)
sx, sy, sw_, sh_ = 0.7, 2.6, 2.15, 1.9
gap = 0.28
for i, (title, body) in enumerate(steps):
    cx = sx + i * (sw_ + gap)
    col = EMBER if i == 3 else STEEL
    rect(s, cx, sy, sw_, sh_, fill=WHITE, shadow=True, radius=0.08)
    rect(s, cx, sy, sw_, 0.45, fill=col, radius=0.35)
    txt(s, cx, sy, sw_, 0.45, title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx+0.15, sy+0.6, sw_-0.3, sh_-0.7, body, size=11.5, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)
    if i < n - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx+sw_+0.02), Inches(sy+sh_/2-0.14), Inches(0.24), Inches(0.28))
        arr.fill.solid(); arr.fill.fore_color.rgb = EMBER
        arr.line.fill.background(); arr.shadow.inherit = False
txt(s, 0.7, 4.95, 11.5, 0.6,
    "Каждый узел графа несёт источник, дословную цитату и уровень достоверности — модель верификации знаний из ТЗ.",
    size=13.5, italic=True, color=SLATE)
page_num(s, 5)

# ============================================================ SLIDE 6 — DEMO QUERIES
s = add_slide(LIGHT_BG)
kicker(s, "Демонстрация")
txt(s, 0.7, 0.95, 11, 0.7, "4 эталонных запроса из ТЗ — с числами и подтверждено", size=28, bold=True, color=INK, font=FONT_HEAD)

demo_q = [
    ("Обессоливание воды", "Сульфаты/хлориды 200-300 мг/л, сухой остаток ≤1000 мг/дм³", "26 условий · 24 утверждения"),
    ("Циркуляция католита", "Технические решения EW никеля в мировой практике, оптимальная скорость потока", "60 утверждений с числами (л/мин, м³/ч)"),
    ("Au, Ag, МПГ", "Распределение между штейном и шлаком за последние 10 лет", "60 утверждений, коэффициенты распределения"),
    ("Закачка шахтных вод", "Способы в России и за рубежом, технико-экономические показатели", "6 утверждений · гео-фильтр РФ/мир"),
]
cw3, ch3 = 5.6, 1.9
pos3 = [(0.7,2.2),(6.6,2.2),(0.7,4.3),(6.6,4.3)]
for (q, sub, stat), (cx, cy) in zip(demo_q, pos3):
    rect(s, cx, cy, cw3, ch3, fill=WHITE, shadow=True, radius=0.06)
    rect(s, cx, cy, 0.12, ch3, fill=EMBER)
    txt(s, cx+0.35, cy+0.2, cw3-0.6, 0.45, q, size=17, bold=True, color=INK, font=FONT_BODY)
    txt(s, cx+0.35, cy+0.68, cw3-0.6, 0.85, sub, size=12, color=MUTED, line_spacing=1.15)
    txt(s, cx+0.35, cy+ch3-0.5, cw3-0.6, 0.4, stat, size=12.5, bold=True, color=EMBER_DARK)
page_num(s, 6)

# ============================================================ SLIDE 7 — VERIFICATION MODEL
s = add_slide(GRAPHITE)
kicker(s, "Модель верификации", color=EMBER)
txt(s, 0.7, 0.95, 11, 0.7, "Каждый факт — источник, достоверность, дата", size=28, bold=True, color=WHITE, font=FONT_HEAD)

badge_data = [
    ("🟢", "Высокая", "Подтверждено данными / несколькими источниками", GREEN),
    ("🟡", "Средняя", "Единичное упоминание в надёжном источнике", RGBColor(0xE0,0xB0,0x45)),
    ("🔴", "Низкая", "Требует дополнительной проверки", RGBColor(0xD9,0x5B,0x5B)),
]
bx = 0.7
for label, title, body, col in badge_data:
    rect(s, bx, 2.3, 3.75, 1.7, fill=RGBColor(0x25,0x2B,0x35), radius=0.08)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx+0.25), Inches(2.55), Inches(0.4), Inches(0.4))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background(); c.shadow.inherit = False
    txt(s, bx+0.85, 2.5, 2.6, 0.5, title, size=16, bold=True, color=WHITE, font=FONT_BODY)
    txt(s, bx+0.3, 3.1, 3.2, 0.8, body, size=11.5, color=STEEL, line_spacing=1.15)
    bx += 4.05

txt(s, 0.7, 4.35, 11.5, 0.4, "Плюс: статус (active / unverified_quote / superseded), дата актуализации, версионирование фактов",
    size=13, italic=True, color=EMBER)
txt(s, 0.7, 5.0, 11.5, 1.5,
    "233 из 233 извлечённых LLM утверждений прошли string-match верификацию цитаты к исходному тексту документа.\n"
    "Ни одна цитата не была принята без дословного совпадения — расхождения автоматически понижаются до низкой достоверности.",
    size=15, color=WHITE, line_spacing=1.3)
page_num(s, 7)

# ============================================================ SLIDE 8 — TECH STACK
s = add_slide(LIGHT_BG)
kicker(s, "Технологический стек")
txt(s, 0.7, 0.95, 11, 0.7, "Компоненты решения", size=28, bold=True, color=INK, font=FONT_HEAD)

stack = [
    ("Ингест", "PyMuPDF, python-docx, langdetect"),
    ("Числа", "Regex-грамматика (собственная)"),
    ("Тезаурус", "88 канонов, 400+ синонимов RU/EN"),
    ("LLM-извлечение", "Строгая JSON-схема онтологии"),
    ("Граф", "Kùzu (embedded, Cypher)"),
    ("Поиск", "BM25 (rank_bm25) + фильтры графа"),
    ("LLM-инференс", "Yandex AI Studio → OpenRouter fallback"),
    ("API/UI", "FastAPI + Cytoscape.js"),
]
cw4, ch4 = 2.75, 1.15
cols4 = 4
gx4, gy4 = 0.7, 2.15
for i, (a, b) in enumerate(stack):
    cx = gx4 + (i % cols4) * (cw4 + 0.2)
    cy = gy4 + (i // cols4) * (ch4 + 0.2)
    rect(s, cx, cy, cw4, ch4, fill=CARD_BG, radius=0.1)
    txt(s, cx+0.2, cy+0.15, cw4-0.4, 0.35, a, size=13, bold=True, color=EMBER_DARK, font=FONT_BODY)
    txt(s, cx+0.2, cy+0.52, cw4-0.4, ch4-0.6, b, size=10.5, color=INK, line_spacing=1.1)
page_num(s, 8)

# ============================================================ SLIDE 9 — SCALABILITY & FAIR
s = add_slide(LIGHT_BG)
kicker(s, "Масштабируемость и FAIR")
txt(s, 0.7, 0.95, 11, 0.7, "Готово к росту: новые домены, новые источники", size=28, bold=True, color=INK, font=FONT_HEAD)

left_items = [
    "Новый домен = новые записи тезауруса, онтология не меняется",
    "Kùzu держит миллионы узлов; переход на Neo4j — те же Cypher-запросы",
    "Извлечение инкрементально: новый документ проходит слои независимо",
    "Версионирование: Claim.status / superseded_by / created_at",
]
rect(s, 0.7, 2.1, 5.7, 4.5, fill=WHITE, shadow=True, radius=0.06)
txt(s, 1.0, 2.35, 5.1, 0.4, "Расширяемость", size=17, bold=True, color=STEEL, font=FONT_BODY)
bullets(s, 1.0, 2.85, 5.1, 3.5, left_items, size=13, space_after=12)

fair = [
    ("Findable", "BM25 + сущности + фильтры; у каждого факта — ID узла"),
    ("Accessible", "REST API, экспорт JSON; ролевая модель доступа"),
    ("Interoperable", "Онтология сериализуема в RDF/OWL"),
    ("Reusable", "Источник, цитата, достоверность, дата — у каждого факта"),
]
rect(s, 6.7, 2.1, 5.9, 4.5, fill=GRAPHITE, radius=0.06)
txt(s, 7.0, 2.35, 5.2, 0.4, "Соответствие FAIR", size=17, bold=True, color=EMBER, font=FONT_BODY)
fy = 2.9
for title, body in fair:
    txt(s, 7.0, fy, 1.7, 0.9, title, size=13, bold=True, color=WHITE)
    txt(s, 8.7, fy, 3.7, 0.9, body, size=11.5, color=STEEL, line_spacing=1.15)
    fy += 0.85
page_num(s, 9)

# ============================================================ SLIDE 10 — CLOSING
s = add_slide(GRAPHITE)
for (cx, cy, r, col) in [(1.0, 6.3, 1.3, SLATE), (12.2, 1.0, 1.0, EMBER_DARK)]:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background(); c.shadow.inherit = False
txt(s, 0.9, 1.4, 10.5, 0.9, "Граф — источник истины.", size=40, bold=True, color=WHITE, font=FONT_HEAD)
txt(s, 0.9, 2.25, 10.5, 0.9, "LLM — интерфейс, но не факт.", size=40, bold=True, color=EMBER, font=FONT_HEAD)
txt(s, 0.9, 3.5, 10, 0.6, "Научный клубок готов отвечать на сложные многопараметрические запросы", size=17, color=RGBColor(0xCF,0xD6,0xDC))
status = [
    ("101", "документ обработан"),
    ("20 843", "числовых условия"),
    ("233", "утверждения (100% verified)"),
    ("74", "сущности графа"),
]
sx2 = 0.9
for num, label in status:
    w = max(2.0, 0.18*len(label))
    txt(s, sx2, 4.5, w, 0.7, num, size=32, bold=True, color=EMBER, font=FONT_HEAD)
    txt(s, sx2, 5.25, w, 0.6, label, size=12, color=STEEL, line_spacing=1.1)
    sx2 += w + 0.35
txt(s, 0.9, 6.4, 10.5, 0.6, "Репозиторий · демо · презентация — см. сопроводительные материалы подачи",
    size=13, italic=True, color=STEEL)
page_num(s, 10)

prs.save("presentation/nauchny-klubok.pptx")
print("Saved presentation/nauchny-klubok.pptx,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
