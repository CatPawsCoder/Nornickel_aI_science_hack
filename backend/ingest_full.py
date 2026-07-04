# -*- coding: utf-8 -*-
"""Масштабный ингест полного корпуса (data/raw/full, ~2000 файлов).

Поддержка: PDF (текстовый слой), DOCX, PPTX, XLS/XLSX, DOC (best-effort), TXT.
Сканы (PDF без текста) помечаются needs_ocr=true и добираются отдельно (ingest_ocr.py).

Дописывает в тот же docs.jsonl / txt/, что и ingest.py — поэтому старые 101 остаются.
Инкрементально: пропускает уже обработанные id.
"""
import hashlib
import json
import os
import re
import sys
import traceback

import fitz
from docx import Document as DocxDocument
from langdetect import detect, DetectorFactory

DetectorFactory.seed = 0

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FULL = os.path.join(ROOT, "data", "raw", "full")
CACHE = os.path.join(ROOT, "data", "cache")
TXT = os.path.join(CACHE, "txt")
os.makedirs(TXT, exist_ok=True)

YEAR_RE = re.compile(r"(19[89]\d|20[0-2]\d)")
GEO_FOREIGN = re.compile(
    r"(зарубежн|мировой практик|за рубежом|Австрали|Канад|Чили|Финлянд|Кита[йе]|США|"
    r"Outotec|Glencore|Vale|BHP|Boliden|Sherritt|Jinchuan|Sumitomo)", re.I)
GEO_RU = re.compile(
    r"(Норильск|Кольск|Мончегорск|Надеждинск|Талнах|ГМК|отечествен|Россия|российск|"
    r"Урал|Гипроникель|Цветные\s+металлы)", re.I)


def extract_pdf(path):
    doc = fitz.open(path)
    parts = [page.get_text("text") for page in doc]
    n = len(doc)
    doc.close()
    return "\n".join(parts), n


def extract_docx(path):
    d = DocxDocument(path)
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Лист: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def extract_xls(path):
    import xlrd
    wb = xlrd.open_workbook(path)
    parts = []
    for sh in wb.sheets():
        parts.append(f"# Лист: {sh.name}")
        for r in range(sh.nrows):
            cells = [str(sh.cell_value(r, c)) for c in range(sh.ncols)
                     if str(sh.cell_value(r, c)).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def guess_year(filename, text):
    m = YEAR_RE.findall(filename)
    if m:
        return int(m[-1])
    m = YEAR_RE.findall(text[:3000])
    if m:
        from collections import Counter
        return int(Counter(m).most_common(1)[0][0])
    return None


def guess_lang(text):
    s = text[:2000].strip()
    if not s:
        return "unknown"
    try:
        return detect(s)
    except Exception:
        return "unknown"


def geo_hint(text):
    f, r = len(GEO_FOREIGN.findall(text)), len(GEO_RU.findall(text))
    if f and r:
        return "mixed"
    if f:
        return "foreign"
    if r:
        return "ru"
    return "unknown"


def main():
    out_path = os.path.join(CACHE, "docs.jsonl")
    done = set()
    if os.path.exists(out_path):
        for line in open(out_path, encoding="utf-8"):
            try:
                done.add(json.loads(line)["id"])
            except Exception:
                pass

    files = sorted(os.listdir(FULL))
    ok = skip = fail = ocr = 0
    scan_list = []
    with open(out_path, "a", encoding="utf-8") as out:
        for i, fn in enumerate(files):
            path = os.path.join(FULL, fn)
            ext = fn.rsplit(".", 1)[-1].lower() if "." in fn else ""
            doc_id = "f" + hashlib.md5(fn.encode("utf-8")).hexdigest()[:11]
            if doc_id in done:
                skip += 1
                continue
            try:
                npages = None
                if ext == "pdf":
                    text, npages = extract_pdf(path)
                elif ext in ("docx", "docm"):
                    text = extract_docx(path)
                elif ext == "pptx":
                    text = extract_pptx(path)
                elif ext == "xlsx":
                    text = extract_xlsx(path)
                elif ext == "xls":
                    text = extract_xls(path)
                elif ext == "txt":
                    text = open(path, encoding="utf-8", errors="ignore").read()
                elif ext == "doc":
                    # старый бинарный .doc — пробуем как текст, иначе пропускаем на OCR-этап
                    raw = open(path, "rb").read()
                    text = re.sub(rb"[^\x20-\x7e\xc0-\xff\n]", b" ", raw).decode("cp1251", "ignore")
                    if len(re.sub(r"\s", "", text)) < 200:
                        text = ""
                else:
                    skip += 1
                    continue

                text = re.sub(r"[ \t]+", " ", text)
                text = re.sub(r"\n{3,}", "\n\n", text).strip()

                if ext == "pdf" and len(text) < 100:
                    # скан — пометим для OCR
                    scan_list.append({"id": doc_id, "filename": fn, "npages": npages})
                    ocr += 1
                    continue
                if len(text) < 150:
                    skip += 1
                    continue

                txt_path = os.path.join(TXT, doc_id + ".txt")
                with open(txt_path, "w", encoding="utf-8") as tf:
                    tf.write(text)
                rec = {
                    "id": doc_id, "filename": fn, "ext": ext,
                    "title": re.sub(r"^\d{4}_", "", os.path.splitext(fn)[0])[:120],
                    "year": guess_year(fn, text), "lang": guess_lang(text),
                    "geo_hint": geo_hint(text), "n_chars": len(text),
                    "text_path": os.path.relpath(txt_path, ROOT),
                }
                out.write(json.dumps(rec, ensure_ascii=False) + "\n")
                out.flush()
                ok += 1
            except Exception:
                fail += 1
                if fail <= 20:
                    print(f"FAIL {fn}\n{traceback.format_exc()[:300]}", flush=True)
            if (i + 1) % 200 == 0:
                print(f"  {i+1}/{len(files)}  ok={ok} skip={skip} ocr={ocr} fail={fail}", flush=True)

    with open(os.path.join(CACHE, "needs_ocr.json"), "w", encoding="utf-8") as f:
        json.dump(scan_list, f, ensure_ascii=False)
    print(f"\nDONE ok={ok} skip={skip} needs_ocr={ocr} fail={fail}")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()
